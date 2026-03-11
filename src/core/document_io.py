"""Document I/O: read PDF/Word/Excel/PowerPoint/CSV/text, save to Word/text."""

import csv
import os
import tempfile
from pathlib import Path

# Max pages to render as images (prevent excessive API token use)
MAX_PDF_IMAGE_PAGES = 10


def read_pdf(path: str) -> str:
    """Read PDF with PyPDF2 first, fallback to pdfplumber for better extraction.

    Raises ValueError if no text can be extracted (e.g. scanned image PDF).
    """
    # Try PyPDF2 first (fast)
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                texts.append(text)
        if texts:
            return "\n\n".join(texts)
    except Exception:
        pass

    # Fallback: pdfplumber (handles more PDF types)
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text and text.strip():
                    texts.append(text)
        if texts:
            return "\n\n".join(texts)
    except Exception:
        pass

    raise ValueError("NO_TEXT")


def render_pdf_as_images(path: str) -> list[str]:
    """Render PDF pages as temporary PNG images using pypdfium2.

    Returns list of temp image file paths (up to MAX_PDF_IMAGE_PAGES pages).
    Caller is responsible for cleanup, but temp files auto-delete on app exit.
    """
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(path)
    n_pages = min(len(pdf), MAX_PDF_IMAGE_PAGES)
    image_paths = []

    temp_dir = os.path.join(tempfile.gettempdir(), "4bro_pdf")
    os.makedirs(temp_dir, exist_ok=True)

    for i in range(n_pages):
        page = pdf[i]
        # Render at 2x scale for readability (default 72dpi → 144dpi)
        bitmap = page.render(scale=2)
        pil_image = bitmap.to_pil()
        img_path = os.path.join(temp_dir, f"pdf_page_{i + 1}.png")
        pil_image.save(img_path)
        image_paths.append(img_path)

    pdf.close()
    return image_paths


def read_word(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def read_excel(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            header = f"[시트: {ws.title}]"
            sheets.append(header + "\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets) if sheets else "(빈 엑셀 파일)"


def read_powerpoint(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("\t".join(cells))
        if texts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "(빈 프레젠테이션)"


def read_csv_file(path: str) -> str:
    for encoding in ("utf-8", "cp949", "euc-kr"):
        try:
            with open(path, "r", encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                rows = ["\t".join(row) for row in reader if any(row)]
                return "\n".join(rows)
        except UnicodeDecodeError:
            continue
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows = ["\t".join(row) for row in reader if any(row)]
        return "\n".join(rows)


def _read_text_file(path: str) -> str:
    for encoding in ("utf-8", "cp949", "euc-kr"):
        try:
            with open(path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx", ".ppt",
    ".csv",
    ".txt", ".md", ".json", ".xml", ".html", ".htm",
    ".log", ".ini", ".cfg", ".yaml", ".yml",
}


def read_document(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return read_pdf(path)
    elif ext == ".docx":
        return read_word(path)
    elif ext == ".doc":
        raise ValueError(".doc 형식은 지원하지 않습니다. .docx로 변환 후 첨부해주세요.")
    elif ext == ".xlsx":
        return read_excel(path)
    elif ext == ".xls":
        raise ValueError(".xls 형식은 지원하지 않습니다. .xlsx로 변환 후 첨부해주세요.")
    elif ext == ".pptx":
        return read_powerpoint(path)
    elif ext == ".ppt":
        raise ValueError(".ppt 형식은 지원하지 않습니다. .pptx로 변환 후 첨부해주세요.")
    elif ext == ".csv":
        return read_csv_file(path)
    elif ext in (".txt", ".md", ".json", ".xml", ".html", ".htm",
                  ".log", ".ini", ".cfg", ".yaml", ".yml"):
        return _read_text_file(path)
    else:
        raise ValueError(
            f"지원하지 않는 파일 형식: {ext}\n"
            f"지원 형식: PDF, Word, Excel, PowerPoint, CSV, 텍스트 등"
        )


def save_to_word(text: str, path: str):
    from docx import Document
    doc = Document()
    for line in text.split("\n"):
        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line)
    doc.save(path)
